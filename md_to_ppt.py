from markdown import markdown
from pptx import Presentation

# Read the Markdown file
with open('input.md', 'r') as file:
    md_content = file.readlines()

# Create a PowerPoint presentation
presentation = Presentation()

# Add slides for each line in Markdown
for line in md_content:
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title, content = line.split(':') if ':' in line else (line, '')
    slide.shapes.title.text = title.strip()
    if content:
        slide.placeholders[1].text = content.strip()

presentation.save('output.pptx')
